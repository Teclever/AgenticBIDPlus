"""S6 Channel 2 — local extraction (portal-agnostic).

Turns the raw files a bid's fetch staged ($BIDPLUS_RUNTIME_DIR/<portal>/bids/<pk>/) into:
  • cleaned, English, technical/financial TEXT for text-based docs (PDF text layer, Word,
    Excel) — governance boilerplate and non-English blocks dropped;
  • a list of SCAN / IMAGE files to hand WHOLE to Sonnet (no local OCR — Sonnet reads them
    natively via vision);
  • cheap regex "local fields" (EMD / PBG / total value / key dates) for the score-4 preview.

What we DROP (per the brief): non-English text (CID / Devanagari / regional blocks) and
GOVERNANCE boilerplate — general T&C, integrity pact, arbitration, force majeure, liquidated
damages, performance/penalty/warranty clauses, indemnification, governing law, dispute
resolution. What we KEEP: project / implementation scope, technical specs, BOQ, financials.

Design note: regex fields are computed from the RAW text (an EMD/PBG amount often sits inside
a clause we strip for Sonnet), while the text handed to Channel 3 is the CLEANED text.
This module is the SOLE extractor; the §8b Sonnet module (Channel 3) consumes ExtractionResult.
The cleaning patterns are lifted from the proven HAL Pass-2 cleaner.
"""

from __future__ import annotations

import re
import shutil
import subprocess
import tempfile
from dataclasses import dataclass, field
from pathlib import Path

import bidplus.config as config

# A PDF that yields less text than this is treated as a SCAN (image-only) and kept whole.
_MIN_TEXT_CHARS = 120
_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".gif", ".webp"}
# Modern Office XML formats we extract natively.
_DOCX_EXTS = {".docx"}
_XLSX_EXTS = {".xlsx", ".xlsm"}
_PPTX_EXTS = {".pptx"}
# Legacy OLE binary Office formats — no native Python reader; converted via LibreOffice
# headless when available (the Ubuntu deploy box), else marked 'unsupported'.
_LEGACY_EXTS = {".doc", ".xls", ".ppt"}

_CID_RE = re.compile(r"\(cid:\d+\)")

# Governance / legal boilerplate SECTION headers — strip the section (governs the project,
# not its implementation). Lifted from the HAL Pass-2 cleaner.
_BOILERPLATE_SECTION_PATTERNS = tuple(
    re.compile(p, re.IGNORECASE) for p in (
        r"general\s+terms\s+and\s+conditions",
        r"\bterms\s+and\s+conditions\b",
        r"\bintegrity\s+pact\b",
        r"\bvendor\s+declaration\b",
        r"\barbitration\s+clause\b",
        r"\bforce\s+majeure\b",
        r"\bliquidated\s+damages\b",
        r"performance\s+(?:bank\s+)?guarantee\b",
        r"\bpenalty\s+clause\b",
        r"\bwarranty\s+clause\b",
        r"\bindemnification\b",
        r"\bgoverning\s+law\b",
        r"\bdispute\s+resolution\b",
    )
)
_BOILERPLATE_INLINE_PHRASES = (
    "authorised signatory", "authorized signatory", "signature of vendor",
    "signature of tenderer", "this is a computer generated", "computer-generated document",
    "duly authorized representative", "for and on behalf of", "place and date",
    "witness :", "witness:", "seal of the", "stamp and signature",
)
_MAJOR_SECTION_RE = re.compile(
    r"^\s*(?:\d+[\.\)]\s|\([a-z]\)\s|[A-Z]\.\s|SECTION\s+\d+|"
    r"Annexure\s+[A-Z0-9]|Appendix\s+[A-Z0-9]|Enclosure\s+\d+)",
    re.IGNORECASE,
)

# Financial regexes (from HAL Pass-2, + a PBG keyword).
_AMOUNT_RE = re.compile(
    r"(?:Rs\.?\s*|INR\s*|₹\s*)(\d[\d,\.]*(?:\s*(?:Lakhs?|Lacs?|Crores?|Cr|L))?(?:/-)?)"
    r"|(\d[\d,\.]*\s*(?:Lakhs?|Lacs?|Crores?|Cr|L)(?:/-)?)",
    re.IGNORECASE,
)
_EMD_KEYWORD_RE = re.compile(r"(?:EMD|Earnest\s+Money(?:\s+Deposit)?)", re.IGNORECASE)
_PBG_KEYWORD_RE = re.compile(
    r"(?:Performance\s+(?:Bank\s+)?Guarantee|Performance\s+Security|\bPBG\b|"
    r"Security\s+Deposit|Bid\s+Security)", re.IGNORECASE)
_VALUE_KEYWORD_RE = re.compile(
    r"(?:Estimated\s+(?:Cost|Value|Project\s+Cost)|Total\s+(?:Value|Project\s+Cost|Contract\s+Value)|"
    r"Contract\s+Value|Project\s+Value|Approximate\s+(?:Cost|Value)|Tender\s+Value)",
    re.IGNORECASE)
_DATE_RE = re.compile(
    r"\b(\d{1,2}[/\-\.]\d{1,2}[/\-\.]\d{2,4}|\d{4}[/\-\.]\d{1,2}[/\-\.]\d{1,2}|"
    r"\d{1,2}\s+[A-Za-z]{3,9}\s+\d{2,4})\b")
_DATE_LABELS = {
    "submission": re.compile(r"(?:bid\s+sub|submission|last\s+date|closing|due\s+date)", re.I),
    "pre_bid": re.compile(r"pre[\-\s]?bid", re.I),
    "opening": re.compile(r"opening", re.I),
    "delivery": re.compile(r"(?:delivery|completion|implementation)\s+(?:period|schedule|date|time)", re.I),
}


@dataclass
class ExtractedDoc:
    doc_name: str
    local_path: str
    kind: str           # 'text' | 'scan' | 'image' | 'unsupported' | 'empty'
    text: str | None = None      # cleaned text (kind == 'text')
    raw_chars: int = 0


@dataclass
class ExtractionResult:
    portal: str
    source_pk: str
    docs: list[ExtractedDoc] = field(default_factory=list)
    combined_text: str = ""              # all cleaned text, for Sonnet
    media_paths: list[str] = field(default_factory=list)  # scans/images → Sonnet whole
    local_fields: dict = field(default_factory=dict)      # score-4 regex preview

    @property
    def text_docs(self) -> list[ExtractedDoc]:
        return [d for d in self.docs if d.kind == "text" and d.text]

    @property
    def unsupported_docs(self) -> list[ExtractedDoc]:
        """Files we could neither extract nor hand to Sonnet (legacy binaries without
        LibreOffice, or unknown formats) — Channel 3 should flag these for manual review."""
        return [d for d in self.docs if d.kind == "unsupported"]


# ── per-format raw extraction ───────────────────────────────────────────────────────

def _pdf_text(path: Path) -> str:
    import pdfplumber
    try:
        with pdfplumber.open(str(path)) as pdf:
            return "\n".join((pg.extract_text() or "") for pg in pdf.pages)
    except Exception:
        return ""


def _docx_text(path: Path) -> str:
    try:
        import docx
        d = docx.Document(str(path))
        return "\n".join(p.text for p in d.paragraphs)
    except Exception:
        return ""


def _xlsx_text(path: Path) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        out: list[str] = []
        for ws in wb.worksheets:
            out.append(f"# sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None and str(c).strip()]
                if cells:
                    out.append("\t".join(cells))
        wb.close()
        return "\n".join(out)
    except Exception:
        return ""


def _pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        out: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            out.append(f"# slide {i}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in para.runs).strip()
                        if t:
                            out.append(t)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [c.text for c in row.cells if c.text.strip()]
                        if cells:
                            out.append("\t".join(cells))
        return "\n".join(out)
    except Exception:
        return ""


def _soffice_bin() -> str | None:
    return shutil.which("soffice") or shutil.which("libreoffice")


def _legacy_to_text(path: Path) -> str | None:
    """Best-effort: convert a legacy OLE Office file (.doc/.xls/.ppt) to PDF via LibreOffice
    headless, then extract its text. Returns None when LibreOffice is absent (e.g. the dev
    Mac) or conversion fails — the caller then marks the doc 'unsupported' for manual review.
    Works wherever `soffice` is on PATH (the Ubuntu deploy box: `apt install libreoffice`)."""
    soffice = _soffice_bin()
    if not soffice:
        return None
    try:
        with tempfile.TemporaryDirectory() as td:
            subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf", "--outdir", td, str(path)],
                capture_output=True, timeout=180, check=True)
            pdfs = list(Path(td).glob("*.pdf"))
            return _pdf_text(pdfs[0]) if pdfs else None
    except Exception:
        return None


def _sniff(path: Path) -> str:
    """Coarse type: 'pdf'|'image'|'docx'|'xlsx'|'pptx'|'legacy'|'other', by magic + ext.

    Modern Office (docx/xlsx/pptx) is a ZIP (PK magic); legacy Office (.doc/.xls/.ppt) is an
    OLE compound file (D0 CF 11 E0 magic). We trust the extension for the specific family."""
    ext = path.suffix.lower()
    try:
        head = path.open("rb").read(8)
    except Exception:
        head = b""
    if head[:4] == b"%PDF" or ext == ".pdf":
        return "pdf"
    if ext in _IMAGE_EXTS:
        return "image"
    if ext in _DOCX_EXTS:
        return "docx"
    if ext in _XLSX_EXTS:
        return "xlsx"
    if ext in _PPTX_EXTS:
        return "pptx"
    if ext in _LEGACY_EXTS or head[:4] == b"\xd0\xcf\x11\xe0":  # OLE compound (legacy Office)
        return "legacy"
    if head[:2] == b"PK":  # zip-based office without a known ext
        return {".docx": "docx", ".xlsx": "xlsx", ".xlsm": "xlsx",
                ".pptx": "pptx"}.get(ext, "other")
    return "other"


# ── cleaning (English-only, governance boilerplate dropped) ──────────────────────────

def _looks_nonenglish(line: str) -> bool:
    letters = [c for c in line if c.isalpha()]
    if not letters:
        return False  # numbers / symbols / table rows — keep
    ascii_letters = sum(1 for c in letters if c.isascii())
    return ascii_letters / len(letters) < 0.5


def _remove_boilerplate_lines(lines: list[str]) -> list[str]:
    result: list[str] = []
    in_bp = False
    for line in lines:
        stripped = line.strip()
        lower = stripped.lower()
        is_bp_header = any(p.search(stripped) for p in _BOILERPLATE_SECTION_PATTERNS)
        is_major_hdr = bool(_MAJOR_SECTION_RE.match(stripped)) or (
            stripped.isupper() and 3 < len(stripped) <= 60 and len(stripped.split()) <= 8)
        if is_bp_header:
            in_bp = True
            continue
        if in_bp:
            if is_major_hdr:
                in_bp = False
                result.append(line)
            continue
        if any(ph in lower for ph in _BOILERPLATE_INLINE_PHRASES):
            continue
        result.append(line)
    return result


def clean_text(raw: str) -> str:
    if not raw:
        return ""
    text = _CID_RE.sub("", raw)
    lines = [ln for ln in text.splitlines() if not _looks_nonenglish(ln)]
    lines = _remove_boilerplate_lines(lines)
    cleaned: list[str] = []
    blank_run = 0
    for line in lines:
        if line.strip() == "":
            blank_run += 1
            if blank_run <= 2:
                cleaned.append(line)
        else:
            blank_run = 0
            cleaned.append(line)
    return "\n".join(cleaned).strip()


# ── regex local fields (from RAW text) ──────────────────────────────────────────────

def _amount_near(text: str, keyword_re: re.Pattern, window: int = 350) -> str | None:
    m = keyword_re.search(text)
    if not m:
        return None
    snippet = text[max(0, m.start() - window // 2): m.end() + window]
    am = _AMOUNT_RE.search(snippet)
    return (am.group(1) or am.group(2) or "").strip() if am else None


def _key_dates(text: str) -> dict:
    out: dict = {}
    for label, kw in _DATE_LABELS.items():
        m = kw.search(text)
        if not m:
            continue
        d = _DATE_RE.search(text, m.end(), m.end() + 200)
        if d:
            out[label] = d.group(1)
    return out


def local_fields(raw: str) -> dict:
    return {
        "emd_value": _amount_near(raw, _EMD_KEYWORD_RE),
        "pbg_value": _amount_near(raw, _PBG_KEYWORD_RE),
        "total_value": _amount_near(raw, _VALUE_KEYWORD_RE),
        "key_dates": _key_dates(raw),
    }


# ── entry point ──────────────────────────────────────────────────────────────────────

def extract_dir(portal: str, source_pk: str, write_sidecars: bool = True) -> ExtractionResult:
    """Extract everything in a bid's staging dir. Text docs → cleaned text (+ optional
    `.txt` sidecar); scans/images → media_paths (whole, for Sonnet); regex fields from raw."""
    staging = config.bid_staging_dir(portal, source_pk)
    res = ExtractionResult(portal=portal, source_pk=source_pk)
    if not staging.is_dir():
        return res

    raw_parts: list[str] = []
    clean_parts: list[str] = []
    for p in sorted(staging.iterdir()):
        if not p.is_file() or p.name.startswith(".") or p.suffix.lower() == ".txt":
            continue
        kind = _sniff(p)
        if kind == "image":
            res.docs.append(ExtractedDoc(p.name, str(p), "image"))
            res.media_paths.append(str(p))
            continue
        if kind == "pdf":
            raw = _pdf_text(p)
            if len((raw or "").strip()) < _MIN_TEXT_CHARS:   # image-only PDF → scan
                res.docs.append(ExtractedDoc(p.name, str(p), "scan", raw_chars=len(raw or "")))
                res.media_paths.append(str(p))
                continue
        elif kind == "docx":
            raw = _docx_text(p)
        elif kind == "xlsx":
            raw = _xlsx_text(p)
        elif kind == "pptx":
            raw = _pptx_text(p)
        elif kind == "legacy":
            raw = _legacy_to_text(p)   # LibreOffice → PDF → text, where available
            if not raw:
                res.docs.append(ExtractedDoc(p.name, str(p), "unsupported"))
                continue
        else:
            res.docs.append(ExtractedDoc(p.name, str(p), "unsupported"))
            continue

        raw = raw or ""
        cleaned = clean_text(raw)
        if not cleaned:
            res.docs.append(ExtractedDoc(p.name, str(p), "empty", raw_chars=len(raw)))
            continue
        res.docs.append(ExtractedDoc(p.name, str(p), "text", text=cleaned, raw_chars=len(raw)))
        raw_parts.append(raw)
        clean_parts.append(f"=== {p.name} ===\n{cleaned}")
        if write_sidecars:
            try:
                (staging / f"{p.name}.txt").write_text(cleaned)
            except Exception:
                pass

    res.combined_text = "\n\n".join(clean_parts)
    res.local_fields = local_fields("\n\n".join(raw_parts))
    return res
